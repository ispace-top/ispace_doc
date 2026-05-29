"""附件预览 API 视图。

支持：PDF 流式预览、docx → HTML 转换、TXT/Markdown 文本预览、MP4 视频流。
"""
import logging
import os
import re
import tempfile
from wsgiref.util import FileWrapper

from django.conf import settings
from django.contrib.auth.decorators import login_required
from django.http import (
    FileResponse, HttpResponse, HttpResponseNotFound,
    HttpResponseBadRequest, JsonResponse, StreamingHttpResponse,
)
from django.shortcuts import get_object_or_404
from django.utils.translation import gettext_lazy as _
from django.views.decorators.http import condition

from backend.apps.doc.models import Attachment

logger = logging.getLogger(__name__)

# 视频流分块大小：1MB
VIDEO_CHUNK_SIZE = 1024 * 1024


def _get_attachment_path(attachment: Attachment) -> str:
    """获取附件的本地文件路径。"""
    return os.path.join(settings.MEDIA_ROOT, attachment.file_path.name)


# ================================================================
# PDF 预览
# ================================================================

@login_required
def preview_pdf(request, attachment_id: int):
    """PDF 文件流式预览，支持 Range 请求（分块加载）。"""
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")
    if not attachment.file_name.lower().endswith(".pdf"):
        return HttpResponseBadRequest("仅支持 PDF 文件预览")

    file_size = os.path.getsize(file_path)
    range_header = request.META.get("HTTP_RANGE", "")

    if range_header:
        # Range 请求：支持前端分块加载
        match = re.match(r"bytes=(\d+)-(\d*)", range_header)
        if not match:
            return HttpResponseBadRequest("无效的 Range 请求")

        start = int(match.group(1))
        end = int(match.group(2)) if match.group(2) else min(start + VIDEO_CHUNK_SIZE, file_size - 1)
        end = min(end, file_size - 1)

        length = end - start + 1
        with open(file_path, "rb") as f:
            f.seek(start)
            data = f.read(length)

        resp = HttpResponse(data, content_type="application/pdf", status=206)
        resp["Content-Range"] = f"bytes {start}-{end}/{file_size}"
        resp["Content-Length"] = str(length)
    else:
        resp = FileResponse(
            FileWrapper(open(file_path, "rb"), VIDEO_CHUNK_SIZE),
            content_type="application/pdf",
        )
        resp["Content-Length"] = str(file_size)

    resp["Accept-Ranges"] = "bytes"
    resp["Content-Disposition"] = f'inline; filename="{attachment.file_name}"'
    return resp


# ================================================================
# docx → HTML 预览
# ================================================================

@login_required
def preview_docx(request, attachment_id: int):
    """docx 文件转换为 HTML 预览。

    依赖: pip install mammoth
    若未安装，返回原始文本提取。
    """
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")
    if not attachment.file_name.lower().endswith(".docx"):
        return HttpResponseBadRequest("仅支持 .docx 文件预览")

    try:
        import mammoth

        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(f)
        return JsonResponse({
            "html": result.value,
            "messages": [str(m) for m in result.messages[:10]],
            "file_name": attachment.file_name,
        })
    except ImportError:
        # 降级方案：提取纯文本
        try:
            import zipfile
            from xml.etree import ElementTree

            with zipfile.ZipFile(file_path) as z:
                with z.open("word/document.xml") as xml_file:
                    tree = ElementTree.parse(xml_file)
                    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                    texts = [node.text for node in tree.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t") if node.text]
                    text = "\n".join(texts)
            return JsonResponse({
                "text": text,
                "file_name": attachment.file_name,
                "note": "安装 mammoth 可获取富文本 HTML 预览: pip install mammoth",
            })
        except Exception as e:
            logger.exception("docx 文本提取失败")
            return HttpResponseBadRequest(f"预览失败: {e}")


# ================================================================
# xlsx → JSON 预览
# ================================================================

@login_required
def preview_xlsx(request, attachment_id: int):
    """xlsx 文件预览，返回各 Sheet 的表格数据（JSON 格式）。

    依赖: pip install openpyxl
    """
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")
    ext = os.path.splitext(attachment.file_name.lower())[1]
    if ext not in {".xlsx", ".xls"}:
        return HttpResponseBadRequest("仅支持 .xlsx/.xls 文件预览")

    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets_data = {}

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            max_cols = 0
            row_limit = 5000  # 最多返回 5000 行

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if row_idx >= row_limit:
                    rows.append([f"... 共 {ws.max_row or '?'} 行，仅展示前 {row_limit} 行"])
                    break
                row_data = [_cell_to_value(cell) for cell in row]
                if any(v is not None for v in row_data):
                    rows.append(row_data)
                    max_cols = max(max_cols, len(row_data))

            sheets_data[sheet_name] = {
                "rows": rows,
                "row_count": len(rows),
                "col_count": max_cols,
                "truncated": ws.max_row and ws.max_row > row_limit,
            }

        wb.close()

        return JsonResponse({
            "sheets": sheets_data,
            "sheet_names": wb.sheetnames,
            "file_name": attachment.file_name,
        })
    except ImportError:
        return JsonResponse({
            "error": "未安装 openpyxl，请执行: pip install openpyxl",
            "file_name": attachment.file_name,
        }, status=500)
    except Exception as e:
        logger.exception("xlsx 解析失败")
        return HttpResponseBadRequest(f"xlsx 解析失败: {e}")


def _cell_to_value(cell) -> str | int | float | None:
    """将 openpyxl 单元格值转为 JSON 可序列化的类型。"""
    if cell is None:
        return None
    if isinstance(cell, (int, float, str, bool)):
        if isinstance(cell, float) and cell == int(cell) and abs(cell) < 1e15:
            return int(cell)
        return cell
    from datetime import datetime, date, time
    if isinstance(cell, datetime):
        return cell.isoformat()
    if isinstance(cell, date):
        return cell.isoformat()
    if isinstance(cell, time):
        return cell.isoformat()
    return str(cell)


# ================================================================
# TXT / Markdown 预览
# ================================================================

@login_required
def preview_text(request, attachment_id: int):
    """文本文件预览（TXT / Markdown / 代码文件），自动检测编码。"""
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")

    ext = os.path.splitext(attachment.file_name.lower())[1]
    text_exts = {".txt", ".md", ".markdown", ".py", ".js", ".ts", ".html", ".css",
                 ".json", ".xml", ".yaml", ".yml", ".csv", ".ini", ".cfg", ".conf",
                 ".log", ".java", ".c", ".cpp", ".h", ".go", ".rs", ".rb", ".php", ".sh", ".sql"}

    if ext not in text_exts:
        return HttpResponseBadRequest(f"不支持预览此文件类型: {ext}")

    # 自动检测编码
    text = ""
    encoding = "utf-8"
    for enc in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read(1024 * 1024)  # 最多读 1MB
            encoding = enc
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    # 检测是否为 Markdown
    is_markdown = ext in {".md", ".markdown"}

    return JsonResponse({
        "text": text,
        "encoding": encoding,
        "is_markdown": is_markdown,
        "file_name": attachment.file_name,
        "size": os.path.getsize(file_path),
    })


# ================================================================
# MP4 视频流
# ================================================================

@login_required
def stream_video(request, attachment_id: int):
    """MP4 视频流播放，支持 Range 请求。"""
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")
    if not attachment.file_name.lower().endswith(".mp4"):
        return HttpResponseBadRequest("仅支持 MP4 视频预览")

    file_size = os.path.getsize(file_path)
    range_header = request.META.get("HTTP_RANGE", "")

    if range_header:
        match = re.match(r"bytes=(\d+)-(\d*)", range_header)
        if not match:
            return HttpResponseBadRequest("无效的 Range 请求")

        start = int(match.group(1))
        end = int(match.group(2)) if match.group(2) else min(start + VIDEO_CHUNK_SIZE * 2, file_size - 1)
        end = min(end, file_size - 1)

        length = end - start + 1
        with open(file_path, "rb") as f:
            f.seek(start)
            data = f.read(length)

        resp = HttpResponse(data, content_type="video/mp4", status=206)
        resp["Content-Range"] = f"bytes {start}-{end}/{file_size}"
        resp["Content-Length"] = str(length)
    else:
        resp = StreamingHttpResponse(
            _file_iterator(file_path, VIDEO_CHUNK_SIZE),
            content_type="video/mp4",
        )
        resp["Content-Length"] = str(file_size)

    resp["Accept-Ranges"] = "bytes"
    resp["Content-Disposition"] = f'inline; filename="{attachment.file_name}"'
    return resp


def _file_iterator(file_path: str, chunk_size: int):
    """文件流迭代器。"""
    with open(file_path, "rb") as f:
        while True:
            chunk = f.read(chunk_size)
            if not chunk:
                break
            yield chunk


# ================================================================
# 文件信息
# ================================================================

@login_required
def file_info(request, attachment_id: int):
    """获取附件基本信息（大小、类型、预览能力）。"""
    attachment = get_object_or_404(Attachment, id=attachment_id)
    file_path = _get_attachment_path(attachment)

    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")

    ext = os.path.splitext(attachment.file_name.lower())[1]
    preview_type = _get_preview_type(ext)
    file_size = os.path.getsize(file_path)

    return JsonResponse({
        "id": attachment.id,
        "file_name": attachment.file_name,
        "file_size": file_size,
        "file_size_human": attachment.file_size,
        "preview_type": preview_type,
        "created_at": attachment.create_time.isoformat() if attachment.create_time else "",
    })


def _get_preview_type(ext: str) -> str:
    """根据扩展名判断预览类型。"""
    mapping = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".txt": "text",
        ".md": "markdown",
        ".py": "code",
        ".js": "code",
        ".ts": "code",
        ".html": "code",
        ".css": "code",
        ".json": "code",
        ".xml": "code",
        ".yaml": "code",
        ".csv": "text",
        ".log": "text",
        ".mp4": "video",
        ".pptx": "pptx",
        ".zip": "zip",
    }
    return mapping.get(ext, "unsupported")


# ================================================================
# 4.2.3 pptx → 图片序列预览 API
# ================================================================

@login_required
def preview_pptx(request, attachment_id):
    """PPTX 幻灯片预览 — 返回幻灯片数量和文本内容（4.2.3）。

    GET /api/preview/pptx/<attachment_id>/

    返回每张幻灯片的文本内容、编号和标题。
    """
    attachment = _get_attachment(attachment_id, request.user)
    if attachment is None:
        return HttpResponseNotFound("附件不存在或无权访问")

    if not attachment.file_path.lower().endswith(".pptx"):
        return HttpResponseBadRequest("仅支持 .pptx 格式")

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return JsonResponse({"error": "python-pptx 未安装"}, status=501)

    file_path = os.path.join(settings.MEDIA_ROOT, attachment.file_path)
    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")

    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            slides.append({
                "index": i + 1,
                "title": texts[0] if texts else f"幻灯片 {i + 1}",
                "text_count": len(texts),
                "preview_text": "\n".join(texts[:5]),
            })

        return JsonResponse({
            "attachment_id": attachment_id,
            "total_slides": len(slides),
            "slides": slides,
        })
    except Exception as e:
        logger.exception("pptx 预览失败")
        return JsonResponse({"error": str(e)}, status=500)


# ================================================================
# 4.3.2 ZIP 目录结构解析 API
# ================================================================

@login_required
def preview_zip(request, attachment_id):
    """ZIP 文件目录结构预览（4.3.2）。

    GET /api/preview/zip/<attachment_id>/

    返回 ZIP 内的文件/目录树。
    """
    import zipfile

    attachment = _get_attachment(attachment_id, request.user)
    if attachment is None:
        return HttpResponseNotFound("附件不存在或无权访问")

    file_path = os.path.join(settings.MEDIA_ROOT, attachment.file_path)
    if not os.path.exists(file_path):
        return HttpResponseNotFound("文件不存在")

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            entries = []
            total_size = 0
            for info in zf.infolist():
                entries.append({
                    "name": info.filename,
                    "size": info.file_size,
                    "compressed_size": info.compress_size,
                    "is_dir": info.is_dir(),
                    "modified": info.date_time,
                })
                total_size += info.file_size

            return JsonResponse({
                "attachment_id": attachment_id,
                "total_entries": len(entries),
                "total_size": total_size,
                "entries": entries,
            })
    except zipfile.BadZipFile:
        return JsonResponse({"error": "无效的 ZIP 文件"}, status=400)
    except Exception as e:
        logger.exception("ZIP 预览失败")
        return JsonResponse({"error": str(e)}, status=500)
